"""
Document Processor  –  ingests PDF, DOCX, PPTX, XLSX, TXT
and returns clean, structured chunks with rich metadata.

Excel aggregation improvements:
  - Stores per-column numeric stats (min, max, mean, sum, count) in a
    dedicated "excel_summary" chunk so the agent can answer aggregation
    queries (total sales, average price, etc.) directly from context.
  - Also stores a structured JSON stats block that the LLM can parse.
"""
from __future__ import annotations

import io
import re
import json
import hashlib
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional

from loguru import logger
import pandas as pd

from config import (
    CHUNK_SIZE, CHUNK_OVERLAP, SUPPORTED_EXTENSIONS
)


# ── Data classes ───────────────────────────────────────────────────────────────

@dataclass
class RawDocument:
    """Single page / sheet extracted from a source file."""
    content:  str
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Chunk:
    """Text chunk ready for embedding + indexing."""
    text:     str
    doc_id:   str
    chunk_id: str
    metadata: Dict[str, Any] = field(default_factory=dict)


# ── Extractors ─────────────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> List[RawDocument]:
    try:
        import PyPDF2
        docs = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    docs.append(RawDocument(
                        content=text,
                        metadata={"page": i + 1, "source_type": "pdf", "file": path.name},
                    ))
        logger.info(f"PDF '{path.name}': extracted {len(docs)} pages")
        return docs
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return []


def _extract_docx(path: Path) -> List[RawDocument]:
    try:
        from docx import Document
        doc = Document(path)
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                full_text += "\n" + " | ".join(c.text for c in row.cells)
        return [RawDocument(
            content=full_text,
            metadata={"source_type": "docx", "file": path.name},
        )]
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return []


def _extract_pptx(path: Path) -> List[RawDocument]:
    try:
        from pptx import Presentation
        prs  = Presentation(path)
        docs = []
        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
            if parts:
                docs.append(RawDocument(
                    content="\n".join(parts),
                    metadata={"slide": i + 1, "source_type": "pptx", "file": path.name},
                ))
        logger.info(f"PPTX '{path.name}': extracted {len(docs)} slides")
        return docs
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return []


def _extract_excel(path: Path) -> List[RawDocument]:
    """
    Smart Excel extraction:
    - Generates a rich summary block with per-column stats (min, max, mean, sum, count).
    - Also embeds a compact JSON stats table so LLMs can answer aggregation queries.
    - Creates row-level batched chunks for row-specific lookups.
    """
    try:
        xl   = pd.ExcelFile(path)
        docs = []

        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            df.columns = [str(c).strip() for c in df.columns]
            df = df.dropna(how="all")

            # ── Rich summary block ────────────────────────────────────────────
            summary_lines = [
                f"Sheet: {sheet}",
                f"Total rows: {len(df)}",
                f"Columns ({len(df.columns)}): {', '.join(df.columns.tolist())}",
                "",
                "=== NUMERIC COLUMN STATISTICS ===",
                "(Use these values to answer aggregation questions like totals, averages, counts)",
            ]

            num_cols = df.select_dtypes(include="number").columns.tolist()
            stats_json: Dict[str, Any] = {"sheet": sheet, "row_count": len(df), "columns": {}}

            for col in num_cols[:20]:   # cap at 20 columns
                s = df[col].dropna()
                if len(s) == 0:
                    continue
                col_stats = {
                    "count": int(len(s)),
                    "sum":   round(float(s.sum()),  4),
                    "mean":  round(float(s.mean()), 4),
                    "min":   round(float(s.min()),  4),
                    "max":   round(float(s.max()),  4),
                }
                stats_json["columns"][col] = col_stats
                summary_lines.append(
                    f"  Column '{col}': "
                    f"count={col_stats['count']}, "
                    f"sum={col_stats['sum']}, "
                    f"mean={col_stats['mean']}, "
                    f"min={col_stats['min']}, "
                    f"max={col_stats['max']}"
                )

            # Text stats for non-numeric columns (top values)
            cat_cols = [c for c in df.columns if c not in num_cols]
            for col in cat_cols[:5]:
                vc = df[col].dropna().astype(str).value_counts()
                if len(vc):
                    top = ", ".join(f"'{v}'({n})" for v, n in vc.head(5).items())
                    summary_lines.append(f"  Column '{col}' (text): top values = {top}")

            # Embed compact JSON stats so LLM can reference exact numbers
            summary_lines.append("")
            summary_lines.append("=== JSON STATS (machine-readable) ===")
            summary_lines.append(json.dumps(stats_json, indent=2))

            docs.append(RawDocument(
                content="\n".join(summary_lines),
                metadata={
                    "sheet":         sheet,
                    "source_type":   "excel_summary",
                    "file":          path.name,
                    "stats":         json.dumps(stats_json),      # available in metadata
                    "has_aggregation_data": True,
                },
            ))

            # ── Row-level blocks ──────────────────────────────────────────────
            batch_size = 20
            for start in range(0, len(df), batch_size):
                chunk_df = df.iloc[start: start + batch_size]
                lines = []
                for _, row in chunk_df.iterrows():
                    row_text = " | ".join(
                        f"{col}: {val}"
                        for col, val in row.items()
                        if pd.notna(val) and str(val).strip()
                    )
                    if row_text:
                        lines.append(row_text)
                if lines:
                    docs.append(RawDocument(
                        content="\n".join(lines),
                        metadata={
                            "sheet":       sheet,
                            "rows":        f"{start + 1}-{start + len(chunk_df)}",
                            "source_type": "excel_rows",
                            "file":        path.name,
                        },
                    ))

        logger.info(f"Excel '{path.name}': extracted {len(docs)} blocks")
        return docs
    except Exception as e:
        logger.error(f"Excel extraction failed: {e}")
        return []


def _extract_txt(path: Path) -> List[RawDocument]:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return [RawDocument(
            content=text,
            metadata={"source_type": "txt", "file": path.name},
        )]
    except Exception as e:
        logger.error(f"TXT extraction failed: {e}")
        return []


# ── Chunker ────────────────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    text = re.sub(r"\s+",       " ", text)
    text = re.sub(r"[^\S\r\n]{3,}", " ", text)
    return text.strip()


def _split_into_chunks(
    text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP
) -> List[str]:
    """Sentence-aware sliding-window chunker."""
    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks: List[str]        = []
    current: List[str]       = []
    current_len: int         = 0

    for sent in sentences:
        words = sent.split()
        w_len = len(words)

        if current_len + w_len > size and current:
            chunks.append(" ".join(current))
            overlap_words: List[str] = []
            overlap_count            = 0
            for s in reversed(current):
                sw = s.split()
                if overlap_count + len(sw) <= overlap:
                    overlap_words = sw + overlap_words
                    overlap_count += len(sw)
                else:
                    break
            current     = [" ".join(overlap_words)] if overlap_words else []
            current_len = len(overlap_words)

        current.append(sent)
        current_len += w_len

    if current:
        chunks.append(" ".join(current))

    return [c for c in chunks if c.strip()]


# ── Public API ─────────────────────────────────────────────────────────────────

def process_file(path: Path) -> List[Chunk]:
    """
    Main entry-point.  Returns list of Chunk objects ready for embedding.
    """
    path = Path(path)
    ext  = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}"
        )

    if ext == ".pdf":
        raw_docs = _extract_pdf(path)
    elif ext == ".docx":
        raw_docs = _extract_docx(path)
    elif ext in (".pptx", ".ppt"):
        raw_docs = _extract_pptx(path)
    elif ext in (".xlsx", ".xls"):
        raw_docs = _extract_excel(path)
    elif ext in (".txt", ".csv"):
        raw_docs = _extract_txt(path)
    else:
        raw_docs = []

    file_hash = hashlib.md5(path.read_bytes()).hexdigest()[:8]
    doc_id    = f"{path.stem}_{file_hash}"

    chunks: List[Chunk] = []
    for doc_idx, raw in enumerate(raw_docs):
        cleaned = _clean_text(raw.content)
        if not cleaned:
            continue

        # Excel summary blocks and short docs stay intact; long text gets chunked
        no_split_types = {"excel_summary"}
        if raw.metadata.get("source_type") in no_split_types or len(cleaned.split()) <= CHUNK_SIZE:
            text_chunks = [cleaned]
        else:
            text_chunks = _split_into_chunks(cleaned)

        for c_idx, text in enumerate(text_chunks):
            chunk_id = f"{doc_id}_d{doc_idx}_c{c_idx}"
            meta     = {
                **raw.metadata,
                "doc_id":      doc_id,
                "file_path":   str(path),
                "chunk_index": c_idx,
                "word_count":  len(text.split()),
            }
            chunks.append(Chunk(
                text=text, doc_id=doc_id, chunk_id=chunk_id, metadata=meta,
            ))

    logger.info(f"'{path.name}' → {len(chunks)} chunks")
    return chunks
